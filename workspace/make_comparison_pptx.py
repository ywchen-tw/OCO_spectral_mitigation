"""make_comparison_pptx.py — regenerate the 7-slide DE/XGB/LinReg TCCON
comparison deck (results/model_comparison/tccon_r100_comparison_tables.pptx)
from the per-tree report CSVs, so the deck always matches the current models.

Sources (all from tccon_comparison_report.py runs, r=100 km / ±60 min):
  tccon_metrics_{ak,direct}_r100km.csv   pooled per-slice footprint RMSE
  tccon_comparison_r100km.csv            per-case biases (station-equal table)
  deep_ensemble/de_prof_mix_*/           feature-set-ablation variant trees

Slides: 1 pooled/QF RMSE (AK) · 2 same (direct) · 3 near-cloud land tail (AK)
· 4 same (direct) · 5 station-equal mean |bias| · 6 feature-set ablation (AK,
2026-07-17 lndo01-retrained variants) · 7 takeaways.

Usage:  PYTHONPATH=src python workspace/make_comparison_pptx.py \
            [--out results/model_comparison/tccon_r100_comparison_tables.pptx]
"""
from __future__ import annotations

import argparse
from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

TREES = {
    "Deep Ensemble": Path("results/model_comparison/deep_ensemble/"
                          "de_beta_nll_prof_reg_foldpca_o05l15_m5/atrain"),
    "XGBoost": Path("results/model_comparison/xgb/xgb_prof_foldpca_o05l15"),
    "LinReg (Ridge)": Path("results/model_comparison/linreg/"
                           "linreg_prof_foldpca_o05l15"),
}
MODELS = list(TREES)

POOLED_SLICES = [
    ("Pooled, QF 0+1", "all", "all", "all"),
    ("Pooled, QF=0 (good)", "qf0", "all", "all"),
    ("Pooled, QF=1", "qf1", "all", "all"),
    ("Ocean, QF 0+1", "all", "ocean", "all"),
    ("Land, QF 0+1", "all", "land", "all"),
    ("Land, QF=1", "qf1", "land", "all"),
]
TAIL_SLICES = [
    ("Land ≤10 km of cloud, QF 0+1", "all", "land", "0–10 km"),
    ("Land ≤10 km of cloud, QF=1", "qf1", "land", "0–10 km"),
    ("Land ≥10 km (far), QF 0+1", "all", "land", "≥10 km"),
]


def load_metrics(ref: str) -> dict[str, pd.DataFrame]:
    return {m: pd.read_csv(TREES[m] / f"tccon_metrics_{ref}_r100km.csv")
            for m in MODELS}


def srow(d, qf, surface, cld):
    m = d[(d.qf_group == qf) & (d.surface == surface) & (d.cld_group == cld)]
    assert len(m) == 1, (qf, surface, cld, len(m))
    return m.iloc[0]


def slice_table(slices, ref):
    mets = load_metrics(ref)
    rows = []
    for label, qf, surface, cld in slices:
        rr = {m: srow(mets[m], qf, surface, cld) for m in MODELS}
        ns = {int(r.n_footprints) for r in rr.values()}
        assert len(ns) == 1, f"footprint mismatch {label}: {ns}"
        before = rr[MODELS[0]].rmse_before
        cells = [f"{label}", f"{ns.pop():,}", f"{before:.2f}"]
        for m in MODELS:
            after = rr[m].rmse_after
            cells.append(f"{after:.2f}  ({1 - after / before:.2f})")
        rows.append(cells)
    return rows


ABL_BASE = Path("results/model_comparison/deep_ensemble")
ABL_VARIANTS = ["no_spec", "no_contam", "no_xco2", "no_xco2_and_spec",
                "no_contam_and_xco2"]
ABL_SLICES = [
    ("Pooled, QF 0+1", "all", "all", "all"),
    ("Pooled, QF=0 (good)", "qf0", "all", "all"),
    ("Pooled, QF=1", "qf1", "all", "all"),
    ("Land, QF 0+1", "all", "land", "all"),
    ("Land ≤10 km of cloud, QF 0+1", "all", "land", "0–10 km"),
    ("Land ≤10 km of cloud, QF=1", "qf1", "land", "0–10 km"),
    ("Land ≥10 km (far), QF 0+1", "all", "land", "≥10 km"),
]


def ablation_table():
    """Rows for the feature-set-ablation slide (AK reference)."""
    trees = {"full": TREES["Deep Ensemble"]}
    trees.update({v: ABL_BASE / f"de_prof_mix_{v}" for v in ABL_VARIANTS})
    mets = {m: pd.read_csv(p / "tccon_metrics_ak_r100km.csv")
            for m, p in trees.items()}
    rows = []
    for label, qf, surface, cld in ABL_SLICES:
        rr = {m: srow(mets[m], qf, surface, cld) for m in trees}
        ns = {int(r.n_footprints) for r in rr.values()}
        assert len(ns) == 1, f"footprint mismatch {label}: {ns}"
        cells = [label, f"{ns.pop():,}", f"{rr['full'].rmse_before:.2f}",
                 f"{rr['full'].rmse_after:.2f}"]
        for v in ABL_VARIANTS:
            cells.append(f"{rr[v].rmse_after:.2f}  "
                         f"({rr[v].rmse_after - rr['full'].rmse_after:+.2f})")
        rows.append(cells)
    return rows


def station_equal(ref: str, qf: str):
    """(before, {model: after}) station-equal mean |bias|."""
    b_col, a_col = (("bias_before", "bias_after") if ref == "ak" else
                    ("bias_before_direct", "bias_after_direct"))
    before = None
    afters = {}
    for m in MODELS:
        d = pd.read_csv(TREES[m] / "tccon_comparison_r100km.csv")
        d = d[(d.surface == "all") & (d.qf_group == qf)]
        pb, pa = [], []
        for _, s in d.groupby("site"):
            w = s["n_oco"].to_numpy(float)
            if w.sum() == 0:
                continue
            pb.append(abs(np.average(s[b_col], weights=w)))
            pa.append(abs(np.average(s[a_col], weights=w)))
        before = float(np.mean(pb))
        afters[m] = float(np.mean(pa))
    return before, afters


def add_table_slide(prs, title, subtitle, header, rows, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25),
                                  Inches(12.5), Inches(0.9))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    nrows, ncols = len(rows) + 1, len(header)
    height = Inches(0.4 * nrows)
    shape = slide.shapes.add_table(nrows, ncols, Inches(0.4), Inches(1.35),
                                   Inches(12.5), height)
    table = shape.table
    for j, h in enumerate(header):
        c = table.cell(0, j)
        c.text = h
        c.text_frame.paragraphs[0].font.size = Pt(12)
        c.text_frame.paragraphs[0].font.bold = True
    for i, row in enumerate(rows, start=1):
        best_j = None
        vals = []
        for j, v in enumerate(row):
            if j >= 3:      # model columns: bold the best RMSE
                try:
                    vals.append((float(str(v).split()[0]), j))
                except ValueError:
                    pass
        if vals:
            best_j = min(vals)[1]
        for j, v in enumerate(row):
            c = table.cell(i, j)
            c.text = str(v)
            f = c.text_frame.paragraphs[0].font
            f.size = Pt(12)
            if best_j is not None and j == best_j:
                f.bold = True
    if note:
        nb = slide.shapes.add_textbox(Inches(0.4), Inches(1.5) + height,
                                      Inches(12.5), Inches(1.0))
        ntf = nb.text_frame
        ntf.word_wrap = True
        ntf.text = note
        ntf.paragraphs[0].font.size = Pt(13)


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25),
                                  Inches(12.5), Inches(0.9))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    bb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4),
                                  Inches(12.3), Inches(5.5))
    btf = bb.text_frame
    btf.word_wrap = True
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(15)
        p.space_after = Pt(10)
    return slide


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="results/model_comparison/"
                                     "tccon_r100_comparison_tables.pptx")
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    header = ["Slice", "n footprints", "Before"] + MODELS

    ak_pool = slice_table(POOLED_SLICES, "ak")
    dr_pool = slice_table(POOLED_SLICES, "direct")
    ak_tail = slice_table(TAIL_SLICES, "ak")
    dr_tail = slice_table(TAIL_SLICES, "direct")

    add_table_slide(
        prs, "TCCON validation — RMSE by model (100 km / ±60 min, AK-harmonized)",
        "Same footprint set & before-column (xco2_bc) for all models · RMSE in "
        "ppm, skill = 1 − RMSE_after/RMSE_before in parentheses · fold-PCA "
        "production models (2026-07-15)",
        header, ak_pool,
        note="DE best on every slice. On QF=0 the models nearly tie — good data "
             "barely needs correcting; the gap lives in QF=1.")
    add_table_slide(
        prs, "TCCON validation — RMSE by model (100 km / ±60 min, direct reference)",
        "Direct = raw TCCON window mean (no averaging-kernel harmonization) · "
        "RMSE in ppm (skill)",
        header, dr_pool,
        note="Ranking unchanged vs AK — the model ordering is reference-invariant.")
    add_table_slide(
        prs, "Near-cloud land tail — AK-harmonized reference",
        "TCCON 100 km / ±60 min · RMSE in ppm (skill)",
        header, ak_tail,
        note="Model choice matters almost entirely where quality is already "
             "flagged and cloud is near.")
    add_table_slide(
        prs, "Near-cloud land tail — direct reference",
        "TCCON 100 km / ±60 min · direct = raw TCCON window mean · RMSE in ppm (skill)",
        header, dr_tail,
        note="Same ordering as the AK reference — the model ranking is "
             "reference-invariant on the near-cloud tail.")

    st_rows = []
    for ref in ("direct", "ak"):
        for qf, qlbl in (("all", "all"), ("qf0", "QF=0"), ("qf1", "QF=1")):
            before, afters = station_equal(ref, qf)
            st_rows.append([ref.upper() if ref == "ak" else "Direct", qlbl,
                            f"{before:.2f}"] +
                           [f"{afters[m]:.2f}" for m in MODELS])
    add_table_slide(
        prs, "Mean per-station |bias| to TCCON (station-equal weights)",
        "18 stations, each station's footprint-weighted mean bias → |·| → "
        "averaged (ppm) · direct = raw TCCON window mean, AK = AK-harmonized",
        ["Reference", "QF"] + ["Before"] + [f"{m} after" for m in MODELS],
        st_rows,
        note="DE's RMSE edge comes from reducing per-footprint scatter and the "
             "tail, not the mean station offset.")

    abl = ablation_table()
    add_table_slide(
        prs, "Feature-set ablation — mix DE, AK-harmonized (fold-PCA, lndo01 variants)",
        "Production full vs no_* group-removal retrains (2026-07-17; identical reg/arch/"
        "folds — variants differ ONLY in feature set) · same footprints & before-column "
        "(asserted) · RMSE in ppm (Δ vs full)",
        ["Slice", "n footprints", "Before", "full"] + ABL_VARIANTS,
        abl,
        note="no_spec and no_contam are TCCON-neutral (Δ ≤ 0.04 ppm everywhere) — "
             "spec + contam are parsimony-droppable; the xco2-departure block carries "
             "the correction (no_xco2 +0.79 pooled, +1.28 near-cloud land QF=1). "
             "Full detail: FEATURESET_ABLATION_QF_2026-07-17.md.")

    def g(rows, i, j):
        return rows[i][j].split()[0]
    bullets = [
        f"Pooled TCCON RMSE — AK: DE {g(ak_pool,0,3)} < XGB {g(ak_pool,0,4)} < "
        f"LinReg {g(ak_pool,0,5)} ppm; direct: {g(dr_pool,0,3)} < "
        f"{g(dr_pool,0,4)} < {g(dr_pool,0,5)}. Ranking is reference-invariant.",
        f"The gap lives in QF=1 and near the cloud edge — on QF=0 the three "
        f"nearly tie (AK: {g(ak_pool,1,3)} / {g(ak_pool,1,4)} / {g(ak_pool,1,5)}; "
        f"direct: {g(dr_pool,1,3)} / {g(dr_pool,1,4)} / {g(dr_pool,1,5)}).",
        f"Near-cloud land QF=1 tail fans out on both references: AK "
        f"{g(ak_tail,1,3)} < {g(ak_tail,1,4)} < {g(ak_tail,1,5)}; direct "
        f"{g(dr_tail,1,3)} < {g(dr_tail,1,4)} < {g(dr_tail,1,5)} ppm.",
        "Station-equal mean bias: DE ≈ XGB — DE's advantage is scatter/tail "
        "reduction, not mean offset. LinReg trails throughout.",
        f"Feature groups: no_spec {abl[0][4].split()[0]} ≈ no_contam "
        f"{abl[0][5].split()[0]} ≈ full {abl[0][3]} pooled (spec + contam "
        f"parsimony-droppable); no_xco2 {abl[0][6].split()[0]} — the "
        "xco2-departure block carries the correction.",
        "Models: fold-PCA production retrains (2026-07-15) — "
        "de_beta_nll_prof_reg_foldpca_o05l15_m5 vs xgb/linreg *_prof_foldpca "
        "pooled 5-fold ensembles; identical features, folds, and TCCON chain; "
        "ablation variants retrained with lndo01 (2026-07-17).",
    ]
    add_bullets_slide(prs, "Takeaways — DE > XGB > LinReg, on both references",
                      bullets)

    Path(args.out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"wrote {args.out}")


if __name__ == "__main__":
    main()
